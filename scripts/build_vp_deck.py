#!/usr/bin/env python3
"""Generate the QSR × WeftOS VP briefing deck as an editable .pptx.

Every element is a native PowerPoint object — title/subtitle placeholders
are layout-bound (so "Change Layout" still works), bullets live in text
frames with theme fonts, and tables are real pptx tables (not images).
Colours are kept minimal so the user can re-theme via the Design tab.

Output: .planning/clients/qsr/qsr-vp-briefing.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / ".planning" / "clients" / "qsr" / "qsr-vp-briefing.pptx"

# ---------------------------------------------------------------------------
# Style helpers
# ---------------------------------------------------------------------------

ACCENT = RGBColor(0x2E, 0x5C, 0x8A)   # muted blue — works on light themes
TABLE_HEADER_BG = RGBColor(0x2E, 0x5C, 0x8A)
TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
BODY_FG = RGBColor(0x1E, 0x1E, 0x1E)
MUTED_FG = RGBColor(0x55, 0x55, 0x55)


def set_title(slide, text, size=32, accent=True):
    title = slide.shapes.title
    title.text = text
    p = title.text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    if accent:
        p.font.color.rgb = ACCENT


def add_bullets(text_frame, bullets, size=18, clear_first=True):
    """Write bullets into a text frame, preserving placeholder so layout-swap still works."""
    if clear_first:
        text_frame.clear()
    for i, item in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        # item can be a string or (text, level) tuple
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = text
        p.level = level
        p.font.size = Pt(size if level == 0 else max(size - 2, 12))
        p.font.color.rgb = BODY_FG
        if level == 0:
            p.font.bold = False


def add_text_box(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=BODY_FG, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return box


def add_table(slide, left, top, width, height, rows):
    """rows[0] is header. Returns the native table shape — fully editable."""
    n_rows = len(rows)
    n_cols = len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            if r == 0:
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = TABLE_HEADER_FG
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            else:
                p.font.size = Pt(12)
                p.font.color.rgb = BODY_FG
    return table


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------

prs = Presentation()
# 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================ SLIDE 1 — Title ==============================
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
set_title(slide1, "Ask what-if of the whole enterprise", size=44)
sub = slide1.placeholders[1]
sub.text = "QSR × WeftOS — VP Briefing"
sub.text_frame.paragraphs[0].font.size = Pt(24)
sub.text_frame.paragraphs[0].font.color.rgb = MUTED_FG

# Tagline below the subtitle
add_text_box(
    slide1,
    Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.5),
    "QSR can already answer what happened. "
    "WeftOS lets QSR ask what would happen if.",
    size=20, italic=True, color=MUTED_FG,
)

# ============================ SLIDE 2 — Measurements ==============================
slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # title only
set_title(slide2, "De-risked against representative data at QSR scale")

add_text_box(
    slide2,
    Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.9),
    "We built and ran the end-to-end system against synthetic operational "
    "data shaped and sized to match QSR's footprint — stores, orders, labor, "
    "inventory, promos, org chart, audits. Known ground truth — the system "
    "grades itself. No real QSR data touched yet.",
    size=15, color=MUTED_FG,
)

# Measured capabilities table
add_table(
    slide2,
    Inches(0.5), Inches(2.6), Inches(12.3), Inches(4.2),
    rows=[
        ["Capability", "Measured"],
        ["Streaming ingest", "180,000 events/sec per worker (~4 workers covers QSR peak)"],
        ["Scenario query end-to-end", "< 1 second, multi-region / multi-brand"],
        ["Per-shard semantic search", "175–580 µs — 33–107× faster than exhaustive search"],
        ["Gap-analysis sweep", "4,600 gaps across 500 stores in 0.15 s"],
        ["Semantic-search recall", "≥ 99 % at default settings (measured, not projected)"],
        ["Audit trail", "BLAKE3 hash-chained, tamper-detectable per entry"],
        ["Privacy", "Employee IDs hashed at ingest; automated PII scan"],
    ],
)

# ============================ SLIDE 3 — Capability ==============================
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
set_title(slide3, "Four things dashboards don't do")

body = slide3.placeholders[1].text_frame
add_bullets(body, [
    "Counterfactual reasoning",
    ("Ask \"if X, then Y\" across geography × time × brand × franchisee. "
     "Simulated through a causal graph, returned with confidence bands — "
     "not a brittle point estimate.", 1),
    "Cryptographic provenance",
    ("Every data point, analytical step, and decision surface carries a "
     "hash-chained audit entry. Blockchain-level integrity guarantees "
     "without operating a consensus network.", 1),
    "Governance as code",
    ("SOX boundaries, franchisee data-sharing rules, and PII regimes "
     "(GDPR / PIPEDA / CCPA / LGPD) enforced at write-time. Violations "
     "blocked before they land.", 1),
    "Programmable substrate",
    ("Extensible to cryptographic multi-party workflows: franchisee "
     "settlement, supply-chain attestation, programmable revenue splits. "
     "Adjacent capabilities beyond initial scope.", 1),
], size=17)

# Footer note
add_text_box(
    slide3,
    Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
    "WeftOS does not replace the data lake. The lake stays authoritative for transactions; "
    "WeftOS holds the reasoning layer above it.",
    size=12, italic=True, color=MUTED_FG,
)

# ============================ SLIDE 4 — Size / speed / hardware ==============================
slide4 = prs.slides.add_slide(prs.slide_layouts[5])
set_title(slide4, "Fits on a commodity server. Scales by adding more.")

add_table(
    slide4,
    Inches(0.5), Inches(1.4), Inches(12.3), Inches(4.8),
    rows=[
        ["Dimension", "Measurement"],
        ["Total storage, 5-year retention",
         "~5 GB across ~192 logical shards; hot/warm/cold tiering automatic"],
        ["HNSW index memory (~20M semantic points)",
         "~4–6 GB resident; per-query working set ~1 GB"],
        ["Scenario query end-to-end",
         "< 1 second"],
        ["Per-shard index build",
         "< 30 seconds"],
        ["Full cold-start build (sharded, parallel)",
         "< 1 minute wall-clock"],
        ["Recommended production hardware",
         "Dual-socket server, 64 GB RAM, 1 TB NVMe (overprovisioned). One warm replica."],
        ["Specialised hardware needed",
         "None. No GPU dependency. No accelerators."],
    ],
)

add_text_box(
    slide4,
    Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7),
    "Not a datacenter conversation. A two-rack-unit conversation.",
    size=18, italic=True, color=ACCENT,
)

# ============================ SLIDE 5 — Delivery ==============================
slide5 = prs.slides.add_slide(prs.slide_layouts[5])
set_title(slide5, "13 weeks from kick-off to production handover")

add_text_box(
    slide5,
    Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.8),
    "Specification, ingest, scenario engine, gap analysis, hardening, audit, "
    "and privacy are already built and measured end-to-end on synthetic data. "
    "The hardest technical risks are behind us.",
    size=15, color=MUTED_FG,
)

add_table(
    slide5,
    Inches(0.5), Inches(2.4), Inches(12.3), Inches(2.9),
    rows=[
        ["Phase", "Weeks", "What lands"],
        ["Data-lake integration", "1–3",
         "Streaming read from QSR's lake; operational dashboards on real data"],
        ["Historical replay", "4–6",
         "One year of QSR history imported into staged shards"],
        ["Calibration", "7–9",
         "Causal edge weights tuned against QSR's observed outcomes"],
        ["Scenario acceptance", "10–11",
         "VP-nominated headline scenario passes accuracy threshold"],
        ["Production rollout", "12–13",
         "Dashboards, audit stream, runbook handover to QSR ops"],
    ],
)

# Footer: needs and team
add_text_box(
    slide5,
    Inches(0.5), Inches(5.5), Inches(6.0), Inches(1.7),
    "Team: 2–3 engineers from us (full-time), one data-lake owner from QSR (part-time).",
    size=13, bold=True, color=BODY_FG,
)

add_text_box(
    slide5,
    Inches(6.8), Inches(5.5), Inches(6.0), Inches(1.7),
    "Needed to start:  (1) data-lake read access,  (2) any A/B or holdout "
    "history,  (3) legal sign-off on franchisee boundaries,  (4) a "
    "nominated first scenario — that becomes the acceptance test.",
    size=13, color=BODY_FG,
)

# ============================ SLIDE 6 — WeftOS vs Snowflake ==============================
slide6 = prs.slides.add_slide(prs.slide_layouts[5])
set_title(slide6, "WeftOS vs plain Snowflake — time to answer (log scale, seconds)")

add_text_box(
    slide6,
    Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.8),
    "Snowflake times assume a warmed medium-sized warehouse with typical "
    "query patterns. WeftOS times are measured against synthetic corpus at "
    "QSR scale. Both values are representative, not worst-case.",
    size=13, italic=True, color=MUTED_FG,
)

# Build an editable, Excel-backed bar chart.
chart_data = CategoryChartData()
chart_data.categories = [
    "Aggregate rollup\n(region × quarter)",
    "Pattern match\n(stores like X last July)",
    "Counterfactual simulation\n(\"if Miami misses budget\")",
    "Gap analysis\n(500-store structural scan)",
    "Streaming ingest latency\n(delta to queryable)",
]
# Seconds per operation. Snowflake values reflect typical cold-warm mix;
# WeftOS values are measured in the Phase-0..4 harness.
chart_data.add_series("Plain Snowflake", (3.0,  45.0, 180.0, 30.0, 60.0))
chart_data.add_series("WeftOS",          (0.5,   0.3,   0.8,  0.15,  0.01))

chart_shape = slide6.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.4),
    chart_data,
)
chart = chart_shape.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False

# Data labels on every bar so the actual numbers are visible even at log scale.
for plot in chart.plots:
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = "0.00\"s\""
    dl.font.size = Pt(10)

# Apply log-scale base-10 to the value axis by inserting <c:logBase val="10"/>
# into the scaling block. This is the documented OOXML way to request a log
# axis and keeps the chart fully editable in PowerPoint (user can revert via
# Format Axis → Axis Options → uncheck Logarithmic scale).
val_ax = chart.value_axis._element
scaling = val_ax.find(qn("c:scaling"))
if scaling is not None:
    existing = scaling.find(qn("c:logBase"))
    if existing is None:
        log_base = etree.SubElement(scaling, qn("c:logBase"))
        log_base.set("val", "10")
    else:
        existing.set("val", "10")

# Readable axis fonts
for ax in (chart.value_axis, chart.category_axis):
    ax.tick_labels.font.size = Pt(11)
    ax.tick_labels.font.color.rgb = BODY_FG

add_text_box(
    slide6,
    Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6),
    "Not in this chart: capabilities Snowflake doesn't natively support — "
    "cryptographic audit chain, governance-as-code, counterfactual engine. "
    "Those are covered in Slide 3.",
    size=12, italic=True, color=MUTED_FG,
)

# ============================ Save ==============================
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Wrote {OUT}")
print(f"Size: {OUT.stat().st_size:,} bytes")
print(f"Slides: {len(prs.slides)}")
